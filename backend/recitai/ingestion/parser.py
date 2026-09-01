"""Document → pages of text, with correct page attribution (spec §9 tasks 1–2).

Page attribution is non-negotiable (invariant I1): a citation that names the wrong page
teaches a student to distrust every other citation. The page map is therefore built
first, from the source's own structure, and never inferred afterwards.

Dispatch is by file extension (D-008). `.pptx` is parsed with `python-pptx` rather than
PyMuPDF, which accepts the format without error but flattens every slide onto one page
and discards speaker notes (I-014). An unknown extension fails loudly rather than falling
through to a parser that would accept the file and mis-attribute it.
"""

import hashlib
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Literal

import structlog

log = structlog.get_logger(__name__)

SourceKind = Literal["pdf", "pptx"]

#: Below this, a PDF is assumed to be scanned images rather than text (§9 task 2).
#: Deliberately NOT applied to slide decks: sparse text is normal there, not a failure
#: signal (I-016).
MIN_CHARS_PER_PAGE_PDF = 100


class UnsupportedFormatError(ValueError):
    """The file extension has no parser. Never fall back to a lenient one."""


class ScannedDocumentError(ValueError):
    """A PDF with too little extractable text. OCR is deferred (§9 task 2)."""


@dataclass(frozen=True)
class ParsedPage:
    """One page of a PDF, or one slide of a deck."""

    number: int  # 1-based, as a student would cite it
    heading: str | None
    text: str
    notes: str = ""

    @property
    def content(self) -> str:
        """Body plus speaker notes.

        Notes are included because on a lecture deck they are often the only prose, and
        they carry the misconceptions §4.5's distractors are built from.
        """
        if self.notes:
            return f"{self.text}\n\n[notes] {self.notes}".strip()
        return self.text


@dataclass(frozen=True)
class PageSpan:
    """Maps a character range of the assembled markdown back to a page."""

    start: int
    end: int
    page: int


@dataclass(frozen=True)
class ParsedDocument:
    filename: str
    source_kind: SourceKind
    pages: list[ParsedPage]
    sha256: str
    markdown: str = ""
    page_map: list[PageSpan] = field(default_factory=list)

    @property
    def page_count(self) -> int:
        return len(self.pages)

    def page_for_offset(self, offset: int) -> int:
        """Resolve a character offset in `markdown` to its page number (§9 task 1)."""
        for span in self.page_map:
            if span.start <= offset < span.end:
                return span.page
        if not self.page_map:
            raise ValueError("document has no pages")
        return self.page_map[-1].page


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as fh:
        for block in iter(lambda: fh.read(1 << 20), b""):
            digest.update(block)
    return digest.hexdigest()


def _assemble(pages: list[ParsedPage]) -> tuple[str, list[PageSpan]]:
    """Concatenate pages into one markdown string, recording each page's char range."""
    parts: list[str] = []
    spans: list[PageSpan] = []
    cursor = 0
    for page in pages:
        block = ""
        if page.heading:
            block += f"## {page.heading}\n\n"
        block += page.content
        block += "\n\n"
        parts.append(block)
        spans.append(PageSpan(start=cursor, end=cursor + len(block), page=page.number))
        cursor += len(block)
    return "".join(parts), spans


# --------------------------------------------------------------------------- PPTX ----


#: Shapes whose tops fall within this many EMU are treated as one visual row.
#: 914400 EMU = 1 inch; ~0.12in is about one line of body text on a slide.
_ROW_BAND_EMU = 109728


@dataclass(frozen=True)
class _PositionedText:
    top: int
    left: int
    text: str
    width: int = 0
    height: int = 0

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height


def _pptx_shape_text(shape: Any, title_id: int | None) -> list[_PositionedText]:
    """Text from one shape, with its position, recursing into groups.

    Position matters because slide decks routinely build "tables" out of dozens of
    individually placed text boxes. python-pptx yields shapes in XML (z-order) order,
    which scrambles them — on a real slide in this corpus, the P3 row precedes the P2 row
    and every value is separated from its column header. Text ordered that way cannot
    support a correct question, so reading order is reconstructed from geometry below.
    """
    out: list[_PositionedText] = []
    if getattr(shape, "shape_id", None) == title_id:
        return out

    top = getattr(shape, "top", None)
    left = getattr(shape, "left", None)

    # Groups nest arbitrarily; their text is real content and is easy to lose.
    if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
        for child in shape.shapes:
            out.extend(_pptx_shape_text(child, title_id))
        return out

    if getattr(shape, "has_table", False) and shape.has_table:
        for row_index, row in enumerate(shape.table.rows):
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                out.append(
                    _PositionedText(
                        top=(top or 0) + row_index,
                        left=left or 0,
                        text=" | ".join(cells),
                        width=getattr(shape, "width", 0) or 0,
                    )
                )
        return out

    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text:
            # A shape with no explicit position inherits the layout's; sort it last rather
            # than pretending it sits at the top of the slide.
            out.append(
                _PositionedText(
                    top=top if top is not None else 1 << 40,
                    left=left if left is not None else 0,
                    text=text,
                    width=getattr(shape, "width", 0) or 0,
                    height=getattr(shape, "height", 0) or 0,
                )
            )
    return out


#: A vertical gap this wide separates two layout blocks — two tables side by side, or a
#: diagram beside its caption. 0.4in on a 13.3in slide. Measured, not guessed: the gutter
#: between `EMP` and `EMP × PAY` on the slide that produced I-027 is 0.59in, and separate
#: shapes are the only way a gap arises at all (a bullet list is one text frame, so it
#: cannot be split by this).
_COLUMN_GUTTER_EMU = 365760

#: A horizontal gap this tall separates stacked blocks. 0.4in.
_ROW_GUTTER_EMU = 365760

#: Depth cap for the recursive cut. Real slides never nest this far; the cap only stops a
#: pathological layout from recursing indefinitely.
_MAX_CUT_DEPTH = 8


def _gaps(spans: list[tuple[int, int]]) -> list[tuple[int, int]]:
    """Empty intervals between merged spans, widest first, as (size, midpoint)."""
    if len(spans) < 2:
        return []
    ordered = sorted(spans)
    found: list[tuple[int, int]] = []
    reach = ordered[0][1]
    for start, end in ordered[1:]:
        if start > reach:
            found.append((start - reach, reach + (start - reach) // 2))
        reach = max(reach, end)
    return sorted(found, reverse=True)


def _dominant_gap(spans: list[tuple[int, int]], floor: int) -> tuple[int, int]:
    """The widest gap, but only if it stands out from the others.

    A table's columns are separated by gaps of roughly equal width, so no single one
    dominates and the table is left whole. Two blocks side by side are separated by one
    gutter much wider than anything inside either block. Requiring dominance is what
    distinguishes them; an absolute threshold alone cannot, because a real gutter measured
    on this corpus is only 0.59in.
    """
    found = _gaps(spans)
    if not found:
        return (0, 0)
    widest, at = found[0]
    if widest < floor:
        return (0, 0)
    if len(found) > 1 and widest < 2 * found[1][0]:
        return (0, 0)
    return (widest, at)


def _cut(items: list[_PositionedText], depth: int = 0) -> list[list[_PositionedText]]:
    """Recursive XY-cut: split on the widest whitespace gutter until none remains.

    Without this, two tables printed side by side are banded into shared rows and their
    columns interleave — which produced a factually wrong question about the shape of a
    Cartesian product (plan/ISSUES.md I-027). Reading order is not a single sweep down the
    slide; it is blocks, each read in turn.
    """
    if len(items) < 2 or depth >= _MAX_CUT_DEPTH:
        return [items]

    x_gap, x_at = _dominant_gap([(i.left, i.right) for i in items], _COLUMN_GUTTER_EMU)
    y_gap, y_at = _dominant_gap([(i.top, i.bottom) for i in items], _ROW_GUTTER_EMU)

    # Columns are cut before rows, even when a horizontal gap is wider. On the I-027
    # slide the tallest vertical gap sits between two left-hand tables while the
    # right-hand table spans the full height — cutting on it first would slice that table
    # in half. Reading order is columns, then rows within each column.
    if x_gap:
        left = [i for i in items if i.left < x_at]
        right = [i for i in items if i.left >= x_at]
        if left and right:
            return _cut(left, depth + 1) + _cut(right, depth + 1)

    if y_gap:
        top = [i for i in items if i.top < y_at]
        bottom = [i for i in items if i.top >= y_at]
        if top and bottom:
            return _cut(top, depth + 1) + _cut(bottom, depth + 1)

    return [items]


def _band_rows(items: list[_PositionedText]) -> list[str]:
    """Group one block's text into visual rows, top to bottom then left to right."""
    if not items:
        return []
    ordered = sorted(items, key=lambda i: (i.top, i.left))
    rows: list[list[_PositionedText]] = [[ordered[0]]]
    for item in ordered[1:]:
        if abs(item.top - rows[-1][0].top) <= _ROW_BAND_EMU:
            rows[-1].append(item)
        else:
            rows.append([item])

    lines: list[str] = []
    for row in rows:
        cells = [c.text for c in sorted(row, key=lambda i: i.left)]
        lines.append(cells[0] if len(cells) == 1 else " | ".join(cells))
    return lines


def _reading_order(items: list[_PositionedText]) -> list[str]:
    """Slide text in reading order: blocks first, then rows within each block.

    Cells sharing a row are joined with " | " so a header row and its values stay on one
    line — the association between "BUDGET" and "150000" is the content of the slide, and
    it is lost the moment they land on separate lines.
    """
    if not items:
        return []
    blocks = _cut(items)
    # Blocks come back in cut order, which is already top-to-bottom / left-to-right.
    lines: list[str] = []
    for block in blocks:
        lines.extend(_band_rows(block))
    return lines


def parse_pptx(path: Path) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[ParsedPage] = []

    for index, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        # Capture the id once: `shapes.title` returns a fresh proxy on each access, so an
        # identity comparison inside the loop silently fails and duplicates the title into
        # the body of every slide.
        title_id = title_shape.shape_id if title_shape is not None else None
        heading = None
        if title_shape is not None and title_shape.has_text_frame:
            heading = title_shape.text_frame.text.strip() or None

        positioned: list[_PositionedText] = []
        for shape in slide.shapes:
            positioned.extend(_pptx_shape_text(shape, title_id))
        body = _reading_order(positioned)

        notes = ""
        if slide.has_notes_slide:
            frame = slide.notes_slide.notes_text_frame
            if frame is not None and frame.text:
                notes = frame.text.strip()

        pages.append(
            ParsedPage(number=index, heading=heading, text="\n".join(body).strip(), notes=notes)
        )

    markdown, page_map = _assemble(pages)
    return ParsedDocument(
        filename=path.name,
        source_kind="pptx",
        pages=pages,
        sha256=sha256_file(path),
        markdown=markdown,
        page_map=page_map,
    )


# ---------------------------------------------------------------------------- PDF ----


def parse_pdf(path: Path) -> ParsedDocument:
    import pymupdf4llm

    page_dicts = pymupdf4llm.to_markdown(str(path), page_chunks=True)
    pages: list[ParsedPage] = []
    for index, page in enumerate(page_dicts, start=1):
        text = (page.get("text") or "").strip()
        pages.append(ParsedPage(number=index, heading=None, text=text))

    total_chars = sum(len(p.text) for p in pages)
    if pages and total_chars / len(pages) < MIN_CHARS_PER_PAGE_PDF:
        raise ScannedDocumentError(
            f"{path.name}: {total_chars} characters across {len(pages)} pages "
            f"({total_chars / len(pages):.0f}/page) is below the {MIN_CHARS_PER_PAGE_PDF}/page "
            f"floor — this looks like a scanned document. OCR is not supported."
        )

    markdown, page_map = _assemble(pages)
    return ParsedDocument(
        filename=path.name,
        source_kind="pdf",
        pages=pages,
        sha256=sha256_file(path),
        markdown=markdown,
        page_map=page_map,
    )


# ----------------------------------------------------------------------- dispatch ----

PARSERS = {".pptx": parse_pptx, ".pdf": parse_pdf}


def parse(path: Path) -> ParsedDocument:
    """Parse a document by extension. Raises rather than guessing (D-008)."""
    suffix = path.suffix.lower()
    parser = PARSERS.get(suffix)
    if parser is None:
        raise UnsupportedFormatError(
            f"{path.name}: no parser for '{suffix}'. Supported: {', '.join(sorted(PARSERS))}. "
            f"Legacy .ppt, Keynote and Google Slides must be re-saved as .pptx — never "
            f"exported to PDF, which discards speaker notes and slide structure (D-008)."
        )
    doc = parser(path)
    log.info(
        "parsed",
        filename=doc.filename,
        kind=doc.source_kind,
        pages=doc.page_count,
        chars=len(doc.markdown),
    )
    return doc
