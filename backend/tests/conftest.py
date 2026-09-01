"""Shared fixtures. Builds a small synthetic deck so ingestion tests never depend on
the copyrighted corpus in materials/ (which is gitignored and absent in CI)."""

from pathlib import Path

import pytest

FIXTURES = Path(__file__).parent / "fixtures"
SYNTHETIC_DECK = FIXTURES / "synthetic_deck.pptx"


@pytest.fixture(scope="session")
def synthetic_deck() -> Path:
    """A 4-slide deck with known content on known slides.

    Regenerated rather than committed as a binary so the expected content lives visibly
    in this file, next to the assertions that depend on it.
    """
    from pptx import Presentation

    FIXTURES.mkdir(exist_ok=True)
    prs = Presentation()
    slides = [
        ("Outline", "Fragmentation\nAllocation", ""),
        (
            "Horizontal Fragmentation",
            "Splits a relation by rows using a predicate.",
            "Students confuse this with vertical.",
        ),
        ("Vertical Fragmentation", "Splits a relation by columns, repeating the key.", ""),
        ("Allocation", "Decide which fragment lives at which site.", ""),
    ]
    for title, body, notes in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    prs.save(str(SYNTHETIC_DECK))
    return SYNTHETIC_DECK


@pytest.fixture(autouse=True)
async def _dispose_engine_between_tests():  # type: ignore[no-untyped-def]
    """Give every test a connection pool bound to its own event loop.

    SQLAlchemy's async pool binds to the loop that first used it, and pytest-asyncio runs
    each test in a fresh loop. Without disposal the second DB-touching test fails with
    "Event loop is closed" — which a broad `except` around a connectivity probe reports
    as "database unreachable", silently skipping real coverage.
    """
    yield
    from recitai.db.session import engine

    await engine.dispose()
